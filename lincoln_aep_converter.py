#!/usr/bin/env python3
"""
Generate lincoln_AEP.pptx and push to GitHub
Run this script to create the presentation file
"""

import sys
from io import BytesIO
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed. Install with: pip install python-pptx")
    sys.exit(1)

# Color scheme from HTML
COLORS = {
    'r5': {'bg': RGBColor(100, 116, 139), 'text': RGBColor(148, 163, 184)},
    'r4': {'bg': RGBColor(59, 130, 246), 'text': RGBColor(147, 197, 253)},
    'r3': {'bg': RGBColor(139, 92, 246), 'text': RGBColor(196, 181, 253)},
    'r2': {'bg': RGBColor(6, 182, 212), 'text': RGBColor(103, 232, 249)},
    'r1': {'bg': RGBColor(0, 169, 157), 'text': RGBColor(45, 212, 191)},
}

ACCENT_COLOR = RGBColor(0, 169, 157)
DARK_BG = RGBColor(7, 16, 31)
LIGHT_TEXT = RGBColor(255, 255, 255)
MUTED_TEXT = RGBColor(200, 200, 200)

# Presentation data
ROWS = [
    {
        'row_id': 'r5',
        'persona': 'Charles Johnson • Facilities & Maintenance',
        'pain': 'SharePoint for maintenance — ungoverned, unauditable, no workflow enforcement or policy compliance',
        'cap_num': 'CAP 05',
        'cap_title': 'Enterprise AI & Workflow Governance',
        'cap_product': 'Platform-wide • Graph Studio contributes lineage, audit & access control',
        'cap_tags': ['Every decision traceable', 'Audit trails', 'Policy enforcement', 'Data lineage', 'Layer-level security'],
        'outcome_title': 'Governed, Auditable Workflows',
        'outcome_desc': 'Every maintenance action logged, policy-enforced, and fully auditable — SharePoint replaced with a governed, compliant platform',
    },
    {
        'row_id': 'r4',
        'persona': 'Charles Johnson • Facilities & Maintenance',
        'pain': 'Work orders disconnected from drawings & asset records — manual stitching required every time',
        'cap_num': 'CAP 04',
        'cap_title': 'Agentic Process Orchestration',
        'cap_product': 'App Studio (Mendix) • Hybrid workforce coordination',
        'cap_tags': ['Governed workflow apps', 'Human + agent coordination', 'Write-back via MCP', 'Replaces ad-hoc tools'],
        'outcome_title': 'Work Orders ↔ Drawings ↔ Assets Unified',
        'outcome_desc': 'Mendix orchestrates the full workflow — work orders, drawings, and asset records connected in one automated, traceable flow',
    },
    {
        'row_id': 'r3',
        'persona': 'Charles Johnson • Facilities & Maintenance',
        'pain': 'Automation needed outside Teamcenter — no low-code path, no way to extend scope to facilities & HR data',
        'cap_num': 'CAP 03',
        'cap_title': 'AI Assisted App Development',
        'cap_product': 'App Studio (Mendix) • Low-code agentic applications',
        'cap_tags': ['Low-code agents', 'Enterprise deployment', 'Connects non-TC data', 'No scope limits', 'Weeks to production'],
        'outcome_title': 'Automation Beyond Teamcenter',
        'outcome_desc': 'Low-code Mendix apps connect facilities, maintenance & HR data alongside TC — no scope limits, production-ready in weeks',
    },
    {
        'row_id': 'r2',
        'persona': 'Jeannie Whited • Multi-System Reporting',
        'pain': 'Facility-wide retrieval needs AI-powered reasoning across asset & operational data — no intelligent layer exists today',
        'cap_num': 'CAP 02',
        'cap_title': 'AI / ML Model Development',
        'cap_product': 'AI Studio • Precision AI grounded in the Knowledge Graph',
        'cap_tags': ['Governed ML models', 'Semantic grounding', 'Predictive analytics', 'Explainable AI'],
        'outcome_title': 'AI-Powered Facility Intelligence',
        'outcome_desc': 'ML models grounded in real asset & operational data — predictive maintenance, anomaly detection, explainable outputs traceable to source',
    },
    {
        'row_id': 'r1',
        'persona': 'Jeannie Whited • Multi-System Reporting',
        'pain': 'Manual data aggregation — one ETL per question, weeks of engineering per new cross-domain report. External integration complexity — data moves, context breaks, no unified cross-domain view',
        'cap_num': 'CAP 01',
        'cap_title': 'Institutional Knowledge Graph',
        'cap_product': 'Rapidminer Graph Studio • The data foundation every other capability depends on',
        'cap_tags': ['One query • Any system', 'No ETL per question', 'Data stays in place', 'W3C RDF / SPARQL / OWL', 'MPP • 10s of billions', 'Scales across every domain'],
        'outcome_title': 'One Query • Any System • No ETL',
        'outcome_desc': 'Semantic overlay connects all sources — cross-domain questions answered in seconds, data stays where it lives, expand without re-architecture',
        'foundation': True,
    },
]

def create_presentation():
    """Create PowerPoint presentation from data"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Create title slide
    create_title_slide(prs)
    
    # Create one slide per row
    for row_data in ROWS:
        create_row_slide(prs, row_data)
    
    return prs

def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Logo
    logo_box = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(1), Inches(0.4))
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = ACCENT_COLOR
    logo_box.line.color.rgb = ACCENT_COLOR
    tf = logo_box.text_frame
    tf.text = "SIEMENS"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = LIGHT_TEXT
    
    # Eyebrow
    eyebrow = slide.shapes.add_textbox(Inches(1.8), Inches(0.55), Inches(7), Inches(0.3))
    tf = eyebrow.text_frame
    tf.text = "PLATFORM RESPONSE"
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # Title
    title = slide.shapes.add_textbox(Inches(1.8), Inches(0.95), Inches(7.5), Inches(1.5))
    tf = title.text_frame
    tf.word_wrap = True
    tf.text = "The Agentic Enterprise Platform — From Pain to Outcome"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    p.line_spacing = 1.2

def create_row_slide(prs, row_data):
    """Create a slide for each row"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    row_id = row_data['row_id']
    color = COLORS[row_id]
    
    # Header with capability number and title
    header_left = Inches(0.5)
    header_top = Inches(0.5)
    header_height = Inches(0.8)
    
    cap_num_box = slide.shapes.add_textbox(header_left, header_top, Inches(0.6), header_height)
    tf = cap_num_box.text_frame
    tf.text = row_data['cap_num']
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = color['text']
    
    cap_title_box = slide.shapes.add_textbox(header_left + Inches(0.8), header_top, Inches(8.7), header_height)
    tf = cap_title_box.text_frame
    tf.word_wrap = True
    tf.text = row_data['cap_title']
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color['text']
    
    # Three-column layout
    col_width = Inches(2.9)
    col_height = Inches(5.8)
    col_left = Inches(0.5)
    col_top = Inches(1.5)
    col_gap = Inches(0.25)
    
    # Column 1: Pain Point
    pain_box = slide.shapes.add_shape(1, col_left, col_top, col_width, col_height)
    pain_box.fill.solid()
    pain_box.fill.fore_color.rgb = DARK_BG
    pain_box.line.color.rgb = color['bg']
    pain_box.line.width = Pt(2)
    
    tf = pain_box.text_frame
    tf.word_wrap = True
    tf.margin_bottom = Inches(0.1)
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    
    # Persona
    p = tf.paragraphs[0]
    p.text = row_data['persona']
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = color['text']
    p.space_after = Pt(6)
    
    # Pain text
    p = tf.add_paragraph()
    p.text = row_data['pain']
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_TEXT
    p.line_spacing = 1.3
    
    # Column 2: Capability
    cap_left = col_left + col_width + col_gap
    cap_box = slide.shapes.add_shape(1, cap_left, col_top, col_width, col_height)
    cap_box.fill.solid()
    cap_box.fill.fore_color.rgb = DARK_BG
    cap_box.line.color.rgb = color['bg']
    cap_box.line.width = Pt(2)
    
    tf = cap_box.text_frame
    tf.word_wrap = True
    tf.margin_bottom = Inches(0.1)
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    
    # Capability product
    p = tf.paragraphs[0]
    p.text = row_data['cap_product']
    p.font.size = Pt(8)
    p.font.color.rgb = MUTED_TEXT
    p.space_after = Pt(8)
    
    # Tags
    for tag in row_data['cap_tags']:
        p = tf.add_paragraph()
        p.text = f"• {tag}"
        p.font.size = Pt(8)
        p.font.color.rgb = color['text']
        p.level = 0
        p.space_after = Pt(4)
    
    # Foundation badge for r1
    if row_data.get('foundation'):
        p = tf.add_paragraph()
        p.text = "◆ Rapidminer Graph Studio — Foundation Layer"
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.space_before = Pt(10)
    
    # Column 3: Outcome
    outcome_left = cap_left + col_width + col_gap
    outcome_box = slide.shapes.add_shape(1, outcome_left, col_top, col_width, col_height)
    outcome_box.fill.solid()
    outcome_box.fill.fore_color.rgb = DARK_BG
    outcome_box.line.color.rgb = color['bg']
    outcome_box.line.width = Pt(2)
    
    tf = outcome_box.text_frame
    tf.word_wrap = True
    tf.margin_bottom = Inches(0.1)
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    
    # Outcome title
    p = tf.paragraphs[0]
    p.text = row_data['outcome_title']
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = color['text']
    p.space_after = Pt(8)
    
    # Outcome description
    p = tf.add_paragraph()
    p.text = row_data['outcome_desc']
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED_TEXT
    p.line_spacing = 1.3

def get_pptx_binary():
    """Generate presentation and return binary data"""
    prs = create_presentation()
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io.read()

if __name__ == '__main__':
    prs = create_presentation()
    prs.save('lincoln_AEP.pptx')
    print("✓ Presentation created: lincoln_AEP.pptx")
